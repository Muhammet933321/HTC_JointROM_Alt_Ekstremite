"""Generate sunum_final.pptx from the ProjeVadi template.

Template structure (preserved):
  - 50.8 x 28.57 cm widescreen, 6 slides, all Blank layout.
  - Each slide already has decorative Freeforms, side bands,
    'SUNUM ŞABLONU' / '2.PROJE VADİ PROJE PAZARI' rotated banners,
    a circular logo area top-right (Freeform 16/18) and a slide number.
  - The main white content panel is a Freeform spanning roughly
    L=8.3..50, T=0..25 cm on slides 2-6.

We don't touch the template chrome. We just add two text boxes per slide
on top of the white area:
  - Title TextBox (top of white area)
  - Body TextBox (below title)
For slide 1 (title slide) we add a centered title + author block instead.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

ROOT = Path(r"c:/Muhammet/Unity/GitHub/HTC_JointROM_Alt_Ekstremite/Assets/Docx")
SRC = ROOT / "PosterSunumu" / "ProjeVadi Proje Pazarı Sunum_Şablonu.pptx"
OUT = ROOT / "PosterSunumu" / "sunum_final.pptx"

ACCENT = "CF0A2D"   # template's red
DARK   = "1B2E4A"   # title navy
BODY   = "000000"

# ----- Slide payloads -----
PROJECT_TITLE = "HTC VIVE Ultimate Tracker ve XR Tabanlı Oyunlaştırılmış Alt Ekstremite Hareket Analizi Sistemi"

def H(text, size=28, color=DARK, bold=True, align="left"):
    return {"text": text, "bold": bold, "size": size, "color": color, "align": align}
def B(text, size=18, color=BODY, bold=False, align="left", level=0, italic=False):
    return {"text": text, "bold": bold, "size": size, "color": color, "align": align,
            "level": level, "italic": italic}
def BUL(text, size=18, level=0):
    return {"text": text, "bold": False, "size": size, "color": BODY,
            "align": "left", "level": level, "bullet": True}

# Slide 1 – Title slide. We place a big centered block over the existing white area.
SLIDE1 = {
    "title_block": [
        {"text": "LİSANS BİTİRME PROJESİ", "bold": True, "size": 20, "color": ACCENT, "align": "center"},
        {"text": PROJECT_TITLE, "bold": True, "size": 34, "color": DARK, "align": "center"},
        {"text": "", "size": 14, "align": "center"},
        {"text": "Muhammet Çiğdem  ·  227231041", "bold": True, "size": 22, "color": BODY, "align": "center"},
        {"text": "Danışman: Dr. Öğr. Üyesi Samet Tonyalı", "bold": False, "size": 20, "color": BODY, "align": "center"},
        {"text": "", "size": 14, "align": "center"},
        {"text": "Gümüşhane Üniversitesi  ·  Yazılım Mühendisliği  ·  Mayıs 2026",
         "bold": False, "size": 18, "color": DARK, "align": "center", "italic": True},
    ]
}

# Slides 2-6: title (red strip) + body (white area).
SLIDES = [
    # Slide 2
    {
        "title": "Problem ve Motivasyon",
        "body": [
            BUL("Diz valgusu, tek ayak dengesizliği ve yük kabulü gibi dinamik paternler statik açı ölçümleriyle yeterince açıklanamaz."),
            BUL("Laboratuvar tipi hareket yakalama sistemleri pahalı ve taşınabilirliği sınırlıdır."),
            BUL("Geleneksel gonyometri değerlendirici bağımlıdır ve anlık ölçüm sınırlılıkları taşır."),
            B(""),
            H("Hedef", size=22, color=ACCENT),
            BUL("HTC VIVE Ultimate Tracker + OpenXR ile erişilebilir bir altyapı kurmak,"),
            BUL("Görev tabanlı bir XR akışına dönüştürmek,"),
            BUL("Tanı iddiası taşımayan destekleyici hareket tarama çıktıları üretmek."),
        ],
    },
    # Slide 3
    {
        "title": "Sistem Mimarisi ve Yöntem",
        "body": [
            H("Donanım & Yazılım", size=22, color=ACCENT),
            BUL("HTC VIVE Ultimate Tracker: pelvis, femur, tibia – ihtiyaç hâlinde üst ekstremite."),
            BUL("Unity 6, OpenXR ve VIVE eklentileri; FK + IK çözücüler."),
            B(""),
            H("4 Katmanlı Veri Akışı", size=22, color=ACCENT),
            BUL("OpenXR Cihaz Katmanı  →  Toplama ve Eşleme Katmanı"),
            BUL("Kinematik Çözüm (FK + tam vücut IK)  →  Uygulama (görev motoru, UI, raporlama)"),
            B(""),
            B("Kalibrasyon:  Q_offset = Q_T⁻¹ · Q_B          Çalışma anı:  Q_B(t) = Q_T(t) · Q_offset",
              size=18, color=DARK, italic=True),
        ],
    },
    # Slide 4
    {
        "title": "Görevler ve Üretilen Metrikler",
        "body": [
            BUL("Kontrollü iniş  →  pik diz valgusu · diz fleksiyonu · sway"),
            BUL("Mini squat  →  diz fleksiyonu · valgus eğilimi"),
            BUL("Tek ayak squat (sağ/sol)  →  diz valgusu · fleksiyon · sway"),
            BUL("Modifiye anterior reach (sağ/sol)  →  reach % · stance valgusu · sway"),
            BUL("Tek ayak denge (sağ/sol)  →  pelvis sway RMS · sway hızı"),
            BUL("Gövde eğimi & yerinde yürüme  →  proksimal kontrol, akış gözlemi"),
            B(""),
            B("Reach % = d_reach / L_limb × 100        Sway_RMS = √( 1/N · Σ (x_i − x̄)² )",
              size=18, color=DARK, italic=True),
            B(""),
            B("Sonuçlar Türkçe yorum ve risk alanı bildirimine dönüştürülmektedir.",
              size=18, color=DARK),
        ],
    },
    # Slide 5
    {
        "title": "Bulgular ve Teknik Katkılar",
        "body": [
            H("Bulgular", size=22, color=ACCENT),
            BUL("Tracker verisinden gerçek zamanlı kalça/diz kinematiği üretimi."),
            BUL("Eşzamanlı tam vücut avatar sürüşü ve görsel geri bildirim."),
            BUL("Görev tabanlı, yorumlanabilir metrik akışı."),
            B(""),
            H("Teknik İyileştirmeler", size=22, color=ACCENT),
            BUL("Hint yön düzeltmesi ile anatomik bükülme düzlemi."),
            BUL("Shin-tracker modu ile alt bacak temsilinin doğruluğu."),
            BUL("Her karede bind-pose sıfırlaması ile drift azaltma."),
            BUL("Çoklu tracker ile doğrudan FK + omurgada CCD tabanlı dağıtım."),
            BUL("Otomatik kemik bulma ve editör test sürücüleri → geliştirme hızında artış."),
        ],
    },
    # Slide 6
    {
        "title": "Sınırlılıklar, Gelecek Çalışmalar ve Sonuç",
        "body": [
            H("Sınırlılıklar", size=22, color=ACCENT),
            BUL("Referans sistemle nicel doğrulama henüz yapılmamıştır."),
            BUL("Görevler literatürdeki LESS / Y-Balance ile birebir eşleşmemektedir."),
            BUL("Trunk lean ve yürüme simülasyonu ayrıntılı klinik ölçüt üretmemektedir."),
            B(""),
            H("Gelecek Çalışmalar", size=22, color=ACCENT),
            BUL("Klinik gonyometre / referans sistemle ROM doğrulaması."),
            BUL("Yürüyüş ve trunk metriklerinin eklenmesi, çoklu tracker rol seçimi."),
            BUL("Terapist paneli ve raporlama modülü."),
            B(""),
            B("Sonuç: ROM ölçümü, full-body avatar sürüşü ve oyunlaştırılmış görev akışı tek mimaride birleştirilerek savunulabilir bir lisans bitirme prototipi sunulmuştur.",
              size=18, color=DARK, italic=True),
        ],
    },
]


# ----- Helpers -----
def add_textbox(slide, left_cm, top_cm, w_cm, h_cm, paragraphs, anchor="top"):
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(w_cm), Cm(h_cm))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    # Replace the default first paragraph with our first spec.
    for i, spec in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Clear existing runs
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        run = p.add_run()
        run.text = spec.get("text", "")
        font = run.font
        font.name = "Times New Roman"
        font.size = Pt(spec.get("size", 18))
        font.bold = bool(spec.get("bold"))
        if spec.get("italic"):
            font.italic = True
        color = spec.get("color")
        if color:
            font.color.rgb = RGBColor.from_string(color)
        p.level = spec.get("level", 0)
        align = spec.get("align", "left")
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }.get(align, PP_ALIGN.LEFT)
        # Manual bullet glyph (no bullet style in default textbox).
        if spec.get("bullet"):
            run.text = "•  " + run.text
    return tb


def main():
    prs = Presentation(str(SRC))

    # --- Slide 1: title ---
    s1 = prs.slides[0]
    # The visible title area on the template slide 1 is roughly L≈8.5..50, T≈3..23.
    add_textbox(s1, left_cm=10.5, top_cm=8.5, w_cm=37.0, h_cm=11.5,
                paragraphs=SLIDE1["title_block"], anchor="middle")

    # --- Slides 2..6: title + body ---
    # White content panel ≈ L 8.5 .. 50, T 0 .. 25 cm.
    # Title placed at top inside red strip-friendly band; body fills the rest.
    for idx, payload in enumerate(SLIDES, start=1):
        slide = prs.slides[idx]
        # Title – sits over the upper portion of the white panel, near the red Freeform.
        add_textbox(slide, left_cm=9.5, top_cm=1.3, w_cm=33.0, h_cm=2.4,
                    paragraphs=[H(payload["title"], size=30, color=ACCENT, bold=True, align="left")],
                    anchor="middle")
        # Body
        add_textbox(slide, left_cm=9.5, top_cm=4.5, w_cm=39.5, h_cm=19.0,
                    paragraphs=payload["body"], anchor="top")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
