"""Deep-inspect a .pptx: recurse into groups and dump every text-bearing shape
with position, size, and full paragraph runs (font, size, color, bold).
Output written as UTF-8 so Turkish chars don't break the encoding.
"""
import sys, io
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

# Force UTF-8 stdout
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

def cm(v):
    try:
        return round(Emu(v).cm, 2)
    except Exception:
        return None

def walk_shape(shape, depth=0):
    indent = "  " * depth
    kind = type(shape).__name__
    try:
        pos = f"L={cm(shape.left)} T={cm(shape.top)} W={cm(shape.width)} H={cm(shape.height)}"
    except Exception:
        pos = "no-geom"
    print(f"{indent}- {kind} name='{shape.name}' {pos}")

    # Recurse into groups
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        try:
            for sub in shape.shapes:
                walk_shape(sub, depth + 1)
        except Exception as e:
            print(f"{indent}  (group iter err: {e})")
        return

    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    for pi, para in enumerate(tf.paragraphs):
        text = "".join(r.text for r in para.runs) or para.text
        if not text.strip():
            continue
        first = para.runs[0] if para.runs else None
        sz = first.font.size.pt if first and first.font.size else None
        bold = first.font.bold if first else None
        fn = first.font.name if first else None
        color = None
        try:
            if first and first.font.color and first.font.color.rgb:
                color = str(first.font.color.rgb)
        except Exception:
            pass
        italic = first.font.italic if first else None
        print(f"{indent}    P[{pi}] sz={sz} bold={bold} italic={italic} font={fn} color={color}")
        print(f"{indent}        text={text!r}")


def main(path):
    prs = Presentation(path)
    print(f"=== {Path(path).name}  ({cm(prs.slide_width)} x {cm(prs.slide_height)} cm, {len(prs.slides)} slides) ===")
    for si, slide in enumerate(prs.slides):
        print(f"\n--- Slide[{si}] layout='{slide.slide_layout.name}' shapes={len(slide.shapes)} ---")
        for shape in slide.shapes:
            walk_shape(shape, 0)


if __name__ == "__main__":
    for a in sys.argv[1:]:
        main(a)
