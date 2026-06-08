"""Generate poster_final.pptx from poster_sablon.pptx by replacing text
inside the existing TextBoxes (97-104). Background image, fonts and panel
positions are preserved from the template.
"""
from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(r"c:/Muhammet/Unity/GitHub/HTC_JointROM_Alt_Ekstremite/Assets/Docx")
SRC = ROOT / "poster" / "poster_sablon.pptx"
OUT = ROOT / "poster" / "poster_final.pptx"

# ---------- Content (mirrors _tools/icerik_planlama.md) ----------
PROJECT_TITLE = "HTC VIVE Ultimate Tracker ve XR Tabanlı Oyunlaştırılmış Alt Ekstremite Hareket Analizi Sistemi"

# Each section: list of paragraphs.
# A paragraph is a dict: {text, bold, size_pt, italic, color_hex, align, level}
def H(text, size=10):  # heading
    return {"text": text, "bold": True, "size": size, "color": "1B2E4A", "align": "left"}

def B(text, size=8):  # body
    return {"text": text, "bold": False, "size": size, "color": "000000", "align": "justify"}

def Bul(text, size=8):  # bullet
    return {"text": "• " + text, "bold": False, "size": size, "color": "000000", "align": "left"}

CONTENT = {
    # Header (top strip)
    "97": [
        {"text": PROJECT_TITLE, "bold": True, "size": 11, "color": "8B0000", "align": "left"},
    ],
    "98": [
        {"text": "Yürütücü: Muhammet Çiğdem", "bold": True, "size": 10, "color": "000000", "align": "left"},
        {"text": "Danışman: Samet Tonyalı", "bold": True, "size": 10, "color": "000000", "align": "left"},
        {"text": "Ekip: GMÜ Yazılım Mühendisliği", "bold": True, "size": 10, "color": "000000", "align": "left"},
    ],

    # Panel 1 – left (Özet)
    "99": [
        H("Özet", 11),
        B("Bu çalışmada HTC VIVE Ultimate Tracker ve OpenXR verisini kullanan Unity tabanlı bir XR hareket analizi sistemi geliştirilmiştir. Sistem; kalça ve diz kinematiğini gerçek zamanlı hesaplamakta, tam vücut avatar görselleştirmesi sağlamakta ve oyunlaştırılmış görevler aracılığıyla hareket paternlerini izlemektedir."),
        B("Kontrollü iniş, mini squat, tek ayak denge, tek ayak squat ve modifiye anterior reach görevlerinde valgus eğilimi, fleksiyon, reach yüzdesi ve salınım temelli metrikler üretilmektedir. Yapı, tanı koyan bir sistemden çok destekleyici bir hareket tarama prototipi olarak konumlandırılmıştır."),
        {"text": "Anahtar Kelimeler: Sanal Gerçeklik, Hareket Açıklığı, Alt Ekstremite, HTC VIVE Ultimate Tracker, OpenXR, Oyunlaştırma, Hareket Tarama.",
         "bold": False, "size": 7.5, "color": "1B2E4A", "align": "justify", "italic": True},
    ],

    # Panel 1 – right (Materyal ve Yöntem)
    "100": [
        H("Projenin Materyali ve Yöntemi", 11),
        Bul("Donanım: HTC VIVE Ultimate Tracker (pelvis, femur, tibia), XR başlık ve kontrolcü."),
        Bul("Yazılım: Unity 6, OpenXR ve VIVE eklentileri."),
        Bul("Algoritma: İleri (FK) ve ters (IK) kinematik çözücüler; kalibrasyon bazlı sensör-kemik ofseti."),
        Bul("Mimari: OpenXR Cihaz → Toplama / Eşleme → Kinematik Çözüm → Uygulama (görev motoru + UI)."),
        B("Kalibrasyon ilişkisi: Q_offset = Q_T⁻¹ · Q_B   ;   Çalışma anı: Q_B^t = Q_T^t · Q_offset."),
        B("Kalça ve diz lokal rotasyonları hiyerarşik FK ile parent-child kemik dönüşümleri üzerinden hesaplanmaktadır."),
    ],

    # Panel 2 – left (Amaç ve Hedefler)
    "101": [
        H("Amaç ve Hedefler", 11),
        B("Amaç: Alt ekstremite hareket analizini XR ortamına taşıyarak sensör verisi, avatar görselleştirme, görev akışı ve kullanıcı geri bildirimini tek yazılım mimarisinde birleştirmek."),
        H("Hedefler", 9),
        Bul("Gerçek zamanlı kalça / diz kinematiği üretimi."),
        Bul("Tam vücut avatar üzerinden görsel geri bildirim."),
        Bul("Görev tabanlı tarama akışı (LESS ve Y-Balance kavramsal referans)."),
        Bul("Türkçe yorum ve risk alanı çıktıları."),
        Bul("Modüler ve sürdürülebilir yazılım yapısı."),
        {"text": "Beklenen başarı ölçütleri: düşük gecikmeli açı akışı, kararlı IK temsili, görev başına yorumlanabilir metrik üretimi.",
         "bold": False, "size": 7.5, "color": "1B2E4A", "align": "justify", "italic": True},
    ],

    # Panel 2 – right (Şekil / Grafik / Tablo)
    "104": [
        H("Görev × Metrik Özeti", 11),
        {"text": "Görev          |  Birincil Metrikler            |  Yorumlanan Patern",
         "bold": True, "size": 7.5, "color": "1B2E4A", "align": "left", "mono": True},
        {"text": "Kontrollü iniş |  Pik valgus, fleksiyon, sway   |  Yük kabulü, iniş kalitesi",
         "bold": False, "size": 7.5, "color": "000000", "align": "left", "mono": True},
        {"text": "Mini squat     |  Diz fleksiyonu, valgus        |  Çift ayak yüklenmesi",
         "bold": False, "size": 7.5, "color": "000000", "align": "left", "mono": True},
        {"text": "Tek ayak squat |  Valgus, fleksiyon, sway       |  Tek taraflı kontrol",
         "bold": False, "size": 7.5, "color": "000000", "align": "left", "mono": True},
        {"text": "Mod. ant.reach |  Reach %, stance valgusu       |  Dinamik denge",
         "bold": False, "size": 7.5, "color": "000000", "align": "left", "mono": True},
        {"text": "Tek ayak denge |  Pelvis sway RMS, hız          |  Denge kontrolü",
         "bold": False, "size": 7.5, "color": "000000", "align": "left", "mono": True},
        {"text": "Gövde eğimi    |  Trunk yönelimi                |  Proksimal kontrol",
         "bold": False, "size": 7.5, "color": "000000", "align": "left", "mono": True},
        B("Reach % = d_reach / L_limb × 100   ;   Sway_RMS = √( (1/N) Σ (x_i − x̄)² )"),
    ],

    # Panel 3 – left (Sonuç ve Öneriler)
    "102": [
        H("Sonuç ve Öneriler", 11),
        B("Sistem; tracker verisinden gerçek zamanlı alt ekstremite kinematiği üretebilmekte, eşzamanlı avatar sürüşü ve görev tabanlı metrik üretimi sağlamaktadır."),
        B("Hint kararlılığı, bind-pose sıfırlaması, uzuv oran uyumu ve shin-tracker temsili gibi pratik IK problemlerine yönelik iyileştirmeler belgelenmiştir."),
        B("Mevcut çıktı tanı amaçlı değildir. Klinik geçerlilik için referans sistemle nicel doğrulama, görev eşiklerinin kullanıcı grubuna kalibrasyonu ve uzman görüşüyle desteklenmiş örnek veri toplanması önerilmektedir."),
    ],

    # Panel 3 – right (Kaynaklar)
    "103": [
        H("Kaynaklar", 11),
        {"text": "[1] Norkin & White. Measurement of Joint Motion. F.A. Davis, 2016.",
         "bold": False, "size": 7, "color": "000000", "align": "left"},
        {"text": "[2] Niehorster ve ark. HTC Vive tracking accuracy. i-Perception 8(3), 2017.",
         "bold": False, "size": 7, "color": "000000", "align": "left"},
        {"text": "[3] HTC. VIVE Ultimate Tracker Specifications, 2024.",
         "bold": False, "size": 7, "color": "000000", "align": "left"},
        {"text": "[4] Renström ve ark. Non-contact ACL injuries. Br. J. Sports Med. 42(6), 2008.",
         "bold": False, "size": 7, "color": "000000", "align": "left"},
        {"text": "[5] Hanzlíková ve ark. LESS systematic review. J. Sci. Med. Sport 24(3), 2021.",
         "bold": False, "size": 7, "color": "000000", "align": "left"},
        {"text": "[6] Kenwright. CCD inverse kinematics. J. Graphics Tools 16(1), 2012.",
         "bold": False, "size": 7, "color": "000000", "align": "left"},
    ],
}

# Geometry overrides (cm) – left, top, width, height. Expand text boxes to fill panels.
# Background panels (visual estimation from template image):
#   Row 1 (Özet / Yöntem)       ≈ T 8.2 → 15.7  (H 7.5)
#   Row 2 (Amaç / Şekil-Grafik) ≈ T 15.9 → 24.2 (H 8.3)
#   Row 3 (Sonuç / Kaynaklar)   ≈ T 24.4 → 28.9 (H 4.5)
GEOM = {
    "97":  (0.6,  4.27, 9.5,  1.30),
    "98":  (10.8, 4.11, 9.5,  1.45),
    "99":  (0.85, 8.30, 9.30, 7.40),
    "100": (10.85,8.30, 9.30, 7.40),
    "101": (0.85, 15.95,9.30, 8.10),
    "104": (10.85,15.95,9.30, 8.10),
    "102": (0.85, 24.30,9.30, 4.50),
    "103": (10.85,24.30,9.30, 4.50),
}

# Map shape name suffix → key in CONTENT.
NAME_TO_KEY = {
    "TextBox 97":  "97",  "TextBox 98":  "98",
    "TextBox 99":  "99",  "TextBox 100": "100",
    "TextBox 101": "101", "TextBox 104": "104",
    "TextBox 102": "102", "TextBox 103": "103",
}


def set_paragraph(tf, idx, spec):
    """Set paragraph idx in text_frame to the given spec. Adds paragraphs if needed."""
    while len(tf.paragraphs) <= idx:
        tf.add_paragraph()
    p = tf.paragraphs[idx]
    # Clear runs
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    p.text = ""  # also nukes any stray text
    run = p.add_run()
    run.text = spec["text"]
    font = run.font
    font.name = "Consolas" if spec.get("mono") else "Times New Roman"
    font.size = Pt(spec["size"])
    font.bold = bool(spec.get("bold"))
    if spec.get("italic"):
        font.italic = True
    color = spec.get("color")
    if color:
        font.color.rgb = RGBColor.from_string(color)
    align = spec.get("align", "left")
    if align == "left":
        p.alignment = PP_ALIGN.LEFT
    elif align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    elif align == "justify":
        p.alignment = PP_ALIGN.JUSTIFY


def fill_textbox(shape, paragraphs, geom):
    # Move/resize the shape.
    left, top, w, h = geom
    shape.left = Cm(left)
    shape.top = Cm(top)
    shape.width = Cm(w)
    shape.height = Cm(h)

    tf = shape.text_frame
    tf.word_wrap = True
    # Wipe existing paragraphs by truncating to one empty paragraph
    # (text_frame must keep at least one paragraph element).
    # Easiest: remove all p children then add fresh ones.
    txBody = tf._txBody
    for p in list(txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p')):
        txBody.remove(p)
    # Re-add one empty paragraph so the frame is well-formed before we fill it.
    from pptx.oxml.ns import qn
    import lxml.etree as ET
    p0 = ET.SubElement(txBody, qn('a:p'))
    # Now populate.
    for i, spec in enumerate(paragraphs):
        set_paragraph(tf, i, spec)


def main():
    prs = Presentation(str(SRC))
    slide = prs.slides[0]
    filled = []
    for shape in slide.shapes:
        key = NAME_TO_KEY.get(shape.name)
        if not key:
            continue
        paragraphs = CONTENT[key]
        geom = GEOM[key]
        fill_textbox(shape, paragraphs, geom)
        filled.append(shape.name)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Filled: {filled}")
    print(f"Wrote: {OUT}")


if __name__ == "__main__":
    main()
