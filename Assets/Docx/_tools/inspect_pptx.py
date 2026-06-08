"""Inspect a .pptx template: list slides, layouts, placeholders, text frames, and basic style info.
Usage: python inspect_pptx.py <path-to-pptx>
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu


def emu_to_cm(v):
    return round(Emu(v).cm, 2) if v is not None else None


def inspect(path: Path):
    prs = Presentation(str(path))
    print(f"=== FILE: {path.name} ===")
    print(f"Slide size: {emu_to_cm(prs.slide_width)} x {emu_to_cm(prs.slide_height)} cm")
    print(f"Slide count: {len(prs.slides)}")
    print(f"Slide masters: {len(prs.slide_masters)}")
    for mi, master in enumerate(prs.slide_masters):
        print(f"  Master[{mi}] layouts: {len(master.slide_layouts)}")
        for li, layout in enumerate(master.slide_layouts):
            print(f"    Layout[{li}] name='{layout.name}' placeholders={len(layout.placeholders)}")
            for ph in layout.placeholders:
                print(f"      ph idx={ph.placeholder_format.idx} type={ph.placeholder_format.type} name='{ph.name}'")

    print("\n--- SLIDES ---")
    for si, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        print(f"\nSlide[{si}] layout='{layout_name}' shapes={len(slide.shapes)}")
        for shi, shape in enumerate(slide.shapes):
            kind = type(shape).__name__
            name = shape.name
            try:
                pos = f"L={emu_to_cm(shape.left)} T={emu_to_cm(shape.top)} W={emu_to_cm(shape.width)} H={emu_to_cm(shape.height)} cm"
            except Exception:
                pos = "no-geom"
            is_ph = shape.is_placeholder
            ph_info = ""
            if is_ph:
                pf = shape.placeholder_format
                ph_info = f" PH(idx={pf.idx}, type={pf.type})"
            print(f"  Shape[{shi}] {kind} name='{name}' {pos}{ph_info}")
            if shape.has_text_frame:
                tf = shape.text_frame
                for pi, para in enumerate(tf.paragraphs):
                    text = "".join(run.text for run in para.runs) or para.text
                    if not text.strip():
                        continue
                    first_run = para.runs[0] if para.runs else None
                    sz = first_run.font.size.pt if first_run and first_run.font.size else None
                    bold = first_run.font.bold if first_run else None
                    name_font = first_run.font.name if first_run else None
                    color = None
                    try:
                        if first_run and first_run.font.color and first_run.font.color.rgb:
                            color = str(first_run.font.color.rgb)
                    except Exception:
                        pass
                    print(f"      P[{pi}] lvl={para.level} sz={sz} bold={bold} font={name_font} color={color} text={text!r}")


if __name__ == "__main__":
    for arg in sys.argv[1:]:
        inspect(Path(arg))
