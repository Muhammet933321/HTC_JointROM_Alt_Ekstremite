"""Generate sunum_final.pptx v2 — fits content into the empty white panels of
the ProjeVadi template.

Important insight: the template's section titles ('PROJE ÖZETİ', 'PROJE
TANITIMI', 'EK AÇIKLAMALAR', 'GÖRSEL ANLATIM', 'TEKNİK GÖRSELLER'), the red
panel headers ('GENEL BAKIŞ', 'PROJE ADI', 'GÖRSEL ANALİZİ' …) and the gray
italic hint texts are baked into the Freeform shapes as picture fills. They
cannot be edited from python-pptx. So we DO NOT duplicate titles; we just
drop our text into the empty white interior of each panel.

Coordinates were measured from the rendered PDF and the Freeform geometry
reported by deep_inspect.py.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

ROOT = Path(r"c:/Muhammet/Unity/GitHub/HTC_JointROM_Alt_Ekstremite/Assets/Docx")
SRC = ROOT / "PosterSunumu" / "ProjeVadi Proje Pazarı Sunum_Şablonu.pptx"
OUT = ROOT / "PosterSunumu" / "sunum_final.pptx"
IMG_DIR = ROOT / "PosterSunumu"

# Brand colors (matched to the template chrome).
RED_DARK = "8B0000"
TEXT = "1A1A1A"
MUTED = "4A4A4A"

PROJECT_TITLE = ("HTC VIVE Ultimate Tracker ve XR Tabanlı Oyunlaştırılmış "
                 "Alt Ekstremite Hareket Analizi Sistemi")


# ----------------------- Helpers -----------------------
def add_image_fit(slide, image_path: Path, box_left_cm, box_top_cm,
                  box_w_cm, box_h_cm, mask_box=True):
    """Insert an image centered inside a box, preserving aspect ratio.
    When mask_box=True, a white rectangle covers the entire box first so the
    underlying template hint text ('GÖRSEL 1' etc., baked into the background
    picture) does not leak around the image's letterboxed edges."""
    from PIL import Image
    from pptx.enum.shapes import MSO_SHAPE
    if mask_box:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(box_left_cm), Cm(box_top_cm), Cm(box_w_cm), Cm(box_h_cm))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
        bg.line.fill.background()  # no border
        bg.shadow.inherit = False

    im = Image.open(image_path)
    iw, ih = im.size
    img_aspect = iw / ih
    box_aspect = box_w_cm / box_h_cm
    if img_aspect > box_aspect:
        w = box_w_cm
        h = w / img_aspect
    else:
        h = box_h_cm
        w = h * img_aspect
    left = box_left_cm + (box_w_cm - w) / 2.0
    top = box_top_cm + (box_h_cm - h) / 2.0
    slide.shapes.add_picture(str(image_path), Cm(left), Cm(top), Cm(w), Cm(h))


def add_textbox(slide, left_cm, top_cm, w_cm, h_cm, paragraphs,
                anchor="top", margin_cm=0.30):
    """Add a non-autosizing text box and fill it with paragraph specs."""
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(w_cm), Cm(h_cm))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Cm(margin_cm)
    tf.margin_right = Cm(margin_cm)
    tf.margin_top = Cm(margin_cm * 0.5)
    tf.margin_bottom = Cm(margin_cm * 0.5)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP,
                          "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]

    for i, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        run = p.add_run()
        text = spec.get("text", "")
        if spec.get("bullet"):
            text = "•  " + text
        run.text = text
        f = run.font
        f.name = spec.get("font", "Calibri")
        f.size = Pt(spec.get("size", 16))
        f.bold = bool(spec.get("bold"))
        if spec.get("italic"):
            f.italic = True
        c = spec.get("color", TEXT)
        if c:
            f.color.rgb = RGBColor.from_string(c)
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY
                       }.get(spec.get("align", "left"), PP_ALIGN.LEFT)
        p.level = spec.get("level", 0)
        if spec.get("space_after"):
            from pptx.util import Pt as _Pt
            p.space_after = _Pt(spec["space_after"])
    return tb


def Body(text, size=14, bold=False, color=TEXT, align="justify",
         space_after=4, italic=False):
    return {"text": text, "size": size, "bold": bold, "color": color,
            "align": align, "space_after": space_after, "italic": italic}

def Bul(text, size=14, color=TEXT, space_after=3):
    return {"text": text, "size": size, "color": color, "align": "left",
            "bullet": True, "space_after": space_after}

def Sub(text, size=15, color=RED_DARK):
    """Sub-heading inside a panel (since the panel header is baked in)."""
    return {"text": text, "size": size, "bold": True, "color": color,
            "align": "left", "space_after": 4}


# ----------------------- Content -----------------------
# Slide 2 — PROJE ÖZETİ → GENEL BAKIŞ (single panel)
SLIDE2_BODY = [
    Body("HTC VIVE Ultimate Tracker ve OpenXR verisini kullanan Unity tabanlı bir "
         "XR hareket analizi sistemi geliştirildi. Sistem; alt ekstremite hareket "
         "açıklığını (ROM) gerçek zamanlı izler, tam vücut avatar görselleştirmesi "
         "sağlar ve oyunlaştırılmış görev akışı ile hareket paternlerini değerlendirir.",
         size=16, space_after=10),
    Body("Kontrollü iniş, mini squat, tek ayak squat, modifiye anterior reach ve "
         "tek ayak denge görevlerinde valgus eğilimi, diz fleksiyonu, reach yüzdesi "
         "ve salınım (sway RMS) gibi metrikler hesaplanır. Sonuçlar kullanıcıya "
         "Türkçe yorum ve risk alanı bildirimi olarak sunulur.",
         size=16, space_after=10),
    Body("Mevcut yapı tanı koyan bir araç değil; klinik değerlendirmeyi destekleyici "
         "bir hareket tarama prototipi olarak konumlandırılmıştır.",
         size=16),
]

# Slide 3 — PROJE TANITIMI (3 columns)
SLIDE3_NAME = [
    Sub("Proje Başlığı", size=14),
    Body(PROJECT_TITLE, size=13, bold=True, color=TEXT, align="left", space_after=10),
    Sub("Tür", size=14),
    Body("Lisans Bitirme Projesi", size=13, color=TEXT, align="left", space_after=10),
    Sub("Konu", size=14),
    Body("Sanal gerçeklik destekli alt ekstremite hareket açıklığı (ROM) "
         "ölçümü ve oyunlaştırılmış değerlendirme.",
         size=12, align="left"),
]
SLIDE3_CAT = [
    Sub("Disiplinler", size=14),
    Bul("Sağlık Teknolojileri", size=13),
    Bul("Biyomekanik / Hareket Analizi", size=13),
    Bul("Sanal Gerçeklik (XR)", size=13),
    Bul("Oyunlaştırma", size=13),
    Bul("Yazılım Mühendisliği", size=13),
    Body(" ", size=8),
    Sub("Anahtar Teknolojiler", size=14),
    Bul("HTC VIVE Ultimate Tracker", size=13),
    Bul("OpenXR  ·  Unity 6", size=13),
    Bul("FK / IK çözücüler", size=13),
]
SLIDE3_TEAM = [
    Sub("Proje Yürütücüsü", size=14),
    Body("Muhammet Çiğdem", size=14, bold=True, align="left", space_after=2),
    Body("Öğrenci No: 227231041", size=12, color=MUTED, align="left", space_after=10),
    Sub("Danışman", size=14),
    Body("Dr. Öğr. Üyesi Samet Tonyalı", size=14, bold=True, align="left", space_after=10),
    Sub("Kurum", size=14),
    Body("Gümüşhane Üniversitesi", size=13, bold=True, align="left", space_after=1),
    Body("Mühendislik ve Doğa Bilimleri Fakültesi", size=12, color=MUTED, align="left", space_after=1),
    Body("Yazılım Mühendisliği Bölümü", size=12, color=MUTED, align="left", space_after=8),
    Body("Mayıs 2026  ·  Gümüşhane", size=12, color=MUTED, align="left"),
]

# Slide 4 — EK AÇIKLAMALAR → DETAYLAR VE NOTLAR (single panel, technical)
SLIDE4_BODY = [
    Sub("Sistem Mimarisi (4 Katman)", size=15),
    Bul("OpenXR Cihaz Katmanı  →  başlık, tracker, kontrolcü, el izleme", size=13, space_after=2),
    Bul("Toplama ve Eşleme  →  cihaz eşleme ve kalibrasyon", size=13, space_after=2),
    Bul("Kinematik Çözüm  →  ROM hesabı + tam vücut IK + kemik sürüşü", size=13, space_after=2),
    Bul("Uygulama  →  klinik açı göstergeleri, avatar, görev motoru, rapor, UI", size=13, space_after=10),

    Sub("Donanım & Yazılım", size=15),
    Bul("HTC VIVE Ultimate Tracker (pelvis, femur, tibia) + XR başlık/kontrolcü", size=13, space_after=2),
    Bul("Unity 6  ·  OpenXR ve VIVE eklentileri  ·  FK / IK çözücüler", size=13, space_after=10),

    Sub("Temel Kinematik İlişki", size=15),
    Body("Kalibrasyon:  Q_offset = Q_T⁻¹ · Q_B          "
         "Çalışma anı:  Q_B(t) = Q_T(t) · Q_offset",
         size=13, align="left", color=MUTED, space_after=10),

    Sub("Görev Seti ve Üretilen Metrikler", size=15),
    Bul("Kontrollü iniş — pik diz valgusu, fleksiyon, sway", size=12, space_after=1),
    Bul("Mini squat & tek ayak squat — fleksiyon, valgus eğilimi", size=12, space_after=1),
    Bul("Modifiye anterior reach — reach %, stance valgusu, sway", size=12, space_after=1),
    Bul("Tek ayak denge — pelvis sway RMS, sway hızı", size=12, space_after=1),
    Bul("Gövde eğimi & yerinde yürüme — proksimal kontrol, akış gözlemi", size=12, space_after=10),

    Sub("Teknik Katkılar", size=15),
    Bul("Hint yön düzeltmesi (anatomik bükülme düzlemi)", size=12, space_after=1),
    Bul("Shin-tracker modu  ·  bind-pose sıfırlaması (drift azaltma)", size=12, space_after=1),
    Bul("Çoklu tracker ile doğrudan FK + omurgada CCD tabanlı dağıtım", size=12, space_after=1),
    Bul("Otomatik kemik bulma ve editör test sürücüleri", size=12),
]

# Slide 5 — GÖRSEL ANLATIM (left image placeholder + right analysis)
# We only fill the right "GÖRSEL ANALİZİ" panel. Image placeholder stays empty
# so the user can drop a Unity screenshot later.
SLIDE5_ANALYSIS = [
    Body("Sistem; OpenXR üzerinden gelen tracker pozlarını kalibrasyon ofseti ile "
         "kemik yönelimine çevirir, ardından kalça ve diz lokal rotasyonlarını "
         "hiyerarşik FK ile çıkarır. Aynı veri akışı tam vücut IK ile avatar "
         "sürüşüne aktarılır.",
         size=13, align="justify", space_after=8),
    Body("Görev motoru bu hareket sinyallerini valgus eğilimi, diz fleksiyonu, "
         "reach yüzdesi ve pelvis sway RMS gibi yorumlanabilir metriklere dönüştürür.",
         size=13, align="justify", space_after=12),
    Sub("Yandaki görsel için öneri", size=12),
    Bul("Kullanıcının HMD ve tracker'larla nötr poz aldığı kalibrasyon ekranı.", size=11, space_after=2),
    Bul("Tam vücut avatar üzerinde kalça/diz açı göstergeleri.", size=11, space_after=0),
]

# Slide 6 — TEKNİK GÖRSELLER (3 image placeholders + bottom notes strip)
# Available area is only ~2.7 cm tall and the strip's left edge has the red
# "ALT METİN / NOTLAR" header. Keep this short, single-line per bullet, 10 pt.
SLIDE6_NOTES = [
    {"text": "Görsel 1 — Tracker yerleşimi: bel, femur ve tibia segmentlerine takılan HTC VIVE Ultimate Tracker'lar.",
     "size": 10, "color": TEXT, "align": "left", "bullet": True, "space_after": 1},
    {"text": "Görsel 2 — Unity Editör görünümü: oyun içi sahne, debug göstergeleri ve hierarchy.",
     "size": 10, "color": TEXT, "align": "left", "bullet": True, "space_after": 1},
    {"text": "Görsel 3 — Oturum sonrası kayıt analiz / replay arayüzü.",
     "size": 10, "color": TEXT, "align": "left", "bullet": True, "space_after": 0},
]


# ----------------------- Layout (cm) -----------------------
# Slide is 50.8 x 28.57 cm.
# Content panel inside Freeform 13/15/17 ≈ L 8.5–50, T 0–25.
# Inside this panel: section title at top, red panel header, white interior,
# bottom hint text. We target only the white interior.

# Single big panel (slides 2 & 4): below "GENEL BAKIŞ"/"DETAYLAR VE NOTLAR" bar
# (red bar bottom ≈ T 5.7), above the bottom hint strip (top ≈ T 22.0).
BIG_PANEL = dict(left_cm=10.6, top_cm=6.4, w_cm=37.3, h_cm=14.7)

# Three columns for slide 3.
# Red panel headers end at ≈ T 5.0. Person icons in right column extend to T ≈ 5.6.
# Bottom hint texts start at ≈ T 21.7.
COL_TOP, COL_H = 6.4, 14.6
COL_LEFT  = dict(left_cm=10.4,  top_cm=COL_TOP, w_cm=12.4, h_cm=COL_H)
COL_MID   = dict(left_cm=23.4,  top_cm=COL_TOP, w_cm=12.4, h_cm=COL_H)
COL_RIGHT = dict(left_cm=36.5,  top_cm=COL_TOP, w_cm=12.4, h_cm=COL_H)

# Slide 5: right-side analysis panel
# Red "GÖRSEL ANALİZİ" bar bottom ≈ T 6.7. Bottom italic hint ≈ T 19.5.
SLIDE5_PANEL = dict(left_cm=30.0, top_cm=7.7, w_cm=18.0, h_cm=11.5)

# Slide 6: bottom notes strip
# Red "ALT METİN / NOTLAR" bar sits ONLY in the upper-left of the wide panel
# (approx L=10.5–20.5, T=19.0–20.0). Italic hint is below at T≈20.7. Best fit:
# place text to the RIGHT of the red bar, vertically centered with it.
SLIDE6_PANEL = dict(left_cm=21.0, top_cm=19.3, w_cm=26.5, h_cm=2.5)

# ---- Image placeholder geometry (measured from rendered template PDF) ----
# Slide 5: large dashed "ANA GÖRSEL / ÇİZİM EKLEYİNİZ" box on the left half.
SLIDE5_IMAGE_BOX = dict(box_left_cm=9.0, box_top_cm=5.4, box_w_cm=19.6, box_h_cm=15.4)

# Slide 6: three dashed image boxes — full dashed frame
# (measured T ≈ 5.0 → 18.0 → height ≈ 13.0 cm).
SLIDE6_IMG_TOP, SLIDE6_IMG_W, SLIDE6_IMG_H = 5.0, 11.0, 13.0
SLIDE6_IMG_LEFT  = dict(box_left_cm=8.7,  box_top_cm=SLIDE6_IMG_TOP, box_w_cm=SLIDE6_IMG_W, box_h_cm=SLIDE6_IMG_H)
SLIDE6_IMG_MID   = dict(box_left_cm=20.7, box_top_cm=SLIDE6_IMG_TOP, box_w_cm=SLIDE6_IMG_W, box_h_cm=SLIDE6_IMG_H)
SLIDE6_IMG_RIGHT = dict(box_left_cm=32.7, box_top_cm=SLIDE6_IMG_TOP, box_w_cm=SLIDE6_IMG_W, box_h_cm=SLIDE6_IMG_H)


SOFFICE = r"C:\Program Files\LibreOffice\program\soffice.exe"


def export_pdf(pptx_path: Path) -> Path:
    """Run LibreOffice headless to convert pptx → pdf, into the same folder."""
    import subprocess
    if not Path(SOFFICE).exists():
        print(f"  [skip pdf] LibreOffice not found at {SOFFICE}")
        return None
    proc = subprocess.run(
        [SOFFICE, "--headless", "--convert-to", "pdf",
         "--outdir", str(pptx_path.parent), str(pptx_path)],
        capture_output=True, text=True, timeout=120,
    )
    pdf_path = pptx_path.with_suffix(".pdf")
    if pdf_path.exists():
        return pdf_path
    print("  [pdf export failed]", proc.stderr or proc.stdout)
    return None


def main():
    prs = Presentation(str(SRC))

    # Slide 1 → cover, untouched.

    # Slide 2 — PROJE ÖZETİ
    add_textbox(prs.slides[1], **BIG_PANEL, paragraphs=SLIDE2_BODY, anchor="top")

    # Slide 3 — PROJE TANITIMI (3 columns, all top-aligned for consistent layout)
    add_textbox(prs.slides[2], **COL_LEFT,  paragraphs=SLIDE3_NAME, anchor="top")
    add_textbox(prs.slides[2], **COL_MID,   paragraphs=SLIDE3_CAT,  anchor="top")
    add_textbox(prs.slides[2], **COL_RIGHT, paragraphs=SLIDE3_TEAM, anchor="top")

    # Slide 4 — EK AÇIKLAMALAR
    add_textbox(prs.slides[3], **BIG_PANEL, paragraphs=SLIDE4_BODY, anchor="top")

    # Slide 5 — GÖRSEL ANLATIM: main image on left, analysis text on right.
    add_image_fit(prs.slides[4], IMG_DIR / "OyunİciGoruntu16-9.png", **SLIDE5_IMAGE_BOX)
    add_textbox(prs.slides[4], **SLIDE5_PANEL, paragraphs=SLIDE5_ANALYSIS, anchor="top")

    # Slide 6 — TEKNİK GÖRSELLER: three images on top, notes at bottom.
    add_image_fit(prs.slides[5], IMG_DIR / "TrackerYerlesimleri.jpeg",      **SLIDE6_IMG_LEFT)
    add_image_fit(prs.slides[5], IMG_DIR / "EditorOyunİciGoruntu3-2.png",   **SLIDE6_IMG_MID)
    add_image_fit(prs.slides[5], IMG_DIR / "KayitAnaliz3-2.png",            **SLIDE6_IMG_RIGHT)
    add_textbox(prs.slides[5], **SLIDE6_PANEL, paragraphs=SLIDE6_NOTES, anchor="top")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote: {OUT}  ({len(prs.slides)} slides)")

    pdf = export_pdf(OUT)
    if pdf:
        print(f"Wrote: {pdf}")


if __name__ == "__main__":
    main()
